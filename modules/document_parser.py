# modules/document_parser.py

import fitz  # PyMuPDF
import pptx

class DocumentParser:
    """
    A utility class to extract plain text from PDF and PowerPoint files.
    """
    def parse(self, file_path: str, file_type: str) -> str:
        """
        Parses the given file and returns the extracted text.

        :param file_path: The path to the file to be parsed.
        :param file_type: The type of the file ('pdf' or 'pptx').
        :return: A string containing the extracted text.
        """
        try:
            if file_type == 'pdf':
                return self._parse_pdf(file_path)
            elif file_type == 'pptx':
                return self._parse_pptx(file_path)
            else:
                return "Error: Unsupported file type."
        except Exception as e:
            return f"Error parsing document: {e}"

    def _parse_pdf(self, file_path: str) -> str:
        """Extracts text from a PDF file."""
        text = ""
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        return text

    def _parse_pptx(self, file_path: str) -> str:
        """Extracts text from a PowerPoint (pptx) file."""
        text = ""
        prs = pptx.Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text